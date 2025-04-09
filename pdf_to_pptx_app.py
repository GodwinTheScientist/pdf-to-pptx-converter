import os
import subprocess

# Function to install Poppler on Streamlit Cloud if not installed
def install_poppler():
    try:
        # Check if poppler is installed by looking for 'pdfinfo'
        subprocess.call(['pdfinfo', '--version'])
    except FileNotFoundError:
        # If 'pdfinfo' is not found, install poppler-utils
        print("Poppler not found. Installing poppler-utils...")
        subprocess.call(['apt-get', 'update'])
        subprocess.call(['apt-get', 'install', '-y', 'poppler-utils'])
    else:
        print("Poppler is already installed.")

install_poppler()  # Call the function to install Poppler

# Continue with the rest of your app code
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import io
import streamlit as st
from pathlib import Path
from io import BytesIO
import fitz  # PyMuPDF

IMAGE_DPI = 150

def convert_pdf(pdf_file, mode):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    filename = Path(pdf_file.name).stem

    if mode == "image":
        images = convert_from_bytes(pdf_file.read(), dpi=IMAGE_DPI)
        for idx, img in enumerate(images):
            temp_img = f"temp_{idx}.png"
            img.save(temp_img)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(temp_img)
    elif mode == "text":
        pdf_file.seek(0)
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        for page in doc:
            text = page.get_text()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.placeholders[1].text = text

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# UI
st.set_page_config(page_title="PDF to PPTX", layout="centered")
st.title("📄 PDF to PPTX Converter")

uploaded_files = st.file_uploader("Upload one or more PDF files", type="pdf", accept_multiple_files=True)
mode = st.selectbox("Conversion mode", ["image", "text"])

if uploaded_files and st.button("Convert"):
    for pdf_file in uploaded_files:
        pptx_file = convert_pdf(pdf_file, mode)
        st.download_button(
            label=f"⬇ Download {pdf_file.name.replace('.pdf', '.pptx')}",
            data=pptx_file,
            file_name=pdf_file.name.replace(".pdf", ".pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
